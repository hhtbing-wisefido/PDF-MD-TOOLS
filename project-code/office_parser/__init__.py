"""
📄 Office文档解析器

支持格式：
- Word: .doc, .docx
- PowerPoint: .ppt, .pptx  
- Excel: .xls, .xlsx
- RTF: .rtf（富文本格式）

依赖：
- python-docx: docx文件解析
- python-pptx: pptx文件解析
- openpyxl: xlsx文件解析
- pywin32: 旧格式(.doc, .ppt, .xls)转换（仅Windows）
- striprtf: rtf文件解析
"""

import os
import io
import tempfile
import shutil
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from datetime import datetime

# 尝试导入Office文档处理库
try:
    from docx import Document as DocxDocument
    from docx.shared import Inches
    from docx.oxml.ns import qn
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    DocxDocument = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    Presentation = None
    MSO_SHAPE_TYPE = None

try:
    from openpyxl import load_workbook
    from openpyxl.drawing.image import Image as OpenpyxlImage
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    load_workbook = None

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# RTF解析
try:
    from striprtf.striprtf import rtf_to_text
    HAS_RTF = True
except ImportError:
    HAS_RTF = False
    rtf_to_text = None

# Windows COM自动化（用于旧格式转换）
HAS_WIN32COM = False
win32com = None
pythoncom = None

def _init_win32com():
    """延迟初始化win32com，避免导入时错误"""
    global HAS_WIN32COM, win32com, pythoncom
    if HAS_WIN32COM:
        return True
    try:
        import win32com.client as _win32com
        import pythoncom as _pythoncom
        win32com = _win32com
        pythoncom = _pythoncom
        HAS_WIN32COM = True
        return True
    except ImportError:
        return False


@dataclass
class ExtractedImage:
    """提取的图片"""
    image_data: bytes
    image_ext: str
    index: int
    width: int = 0
    height: int = 0
    
    def get_filename(self, base_name: str) -> str:
        """生成图片文件名"""
        return f"{base_name}_img{self.index}.{self.image_ext}"


@dataclass
class OfficeContent:
    """Office文档内容"""
    file_type: str  # docx, pptx, xlsx, doc, ppt, xls
    title: str = ""
    text_content: List[Dict[str, Any]] = field(default_factory=list)
    images: List[ExtractedImage] = field(default_factory=list)
    tables: List[List[List[str]]] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    total_images: int = 0
    

# ========== DOCX 解析 ==========

def extract_docx_content(
    file_path: Path,
    output_dir: Path,
    images_subdir: str = "images",
    extract_images: bool = True
) -> OfficeContent:
    """
    提取DOCX文件内容
    
    Args:
        file_path: DOCX文件路径
        output_dir: 输出目录
        images_subdir: 图片子目录名
        extract_images: 是否提取图片
    
    Returns:
        OfficeContent: 提取的内容
    """
    if not HAS_DOCX:
        raise ImportError("需要安装 python-docx: pip install python-docx")
    
    content = OfficeContent(file_type="docx")
    doc = DocxDocument(str(file_path))
    
    # 提取元数据
    core_props = doc.core_properties
    content.metadata = {
        "title": core_props.title or "",
        "author": core_props.author or "",
        "subject": core_props.subject or "",
        "created": str(core_props.created) if core_props.created else "",
        "modified": str(core_props.modified) if core_props.modified else "",
    }
    content.title = core_props.title or file_path.stem
    
    # 创建图片目录
    if extract_images:
        images_dir = output_dir / images_subdir
        images_dir.mkdir(parents=True, exist_ok=True)
    
    image_counter = 0
    
    # 提取段落内容
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        
        # 判断段落类型
        block_type = "paragraph"
        if para.style and para.style.name:
            style_name = para.style.name.lower()
            if "heading 1" in style_name or "标题 1" in style_name:
                block_type = "heading1"
            elif "heading 2" in style_name or "标题 2" in style_name:
                block_type = "heading2"
            elif "heading 3" in style_name or "标题 3" in style_name:
                block_type = "heading3"
            elif "list" in style_name or "bullet" in style_name:
                block_type = "list_item"
        
        content.text_content.append({
            "type": block_type,
            "content": para.text.strip()
        })
    
    # 提取表格
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            table_data.append(row_data)
        content.tables.append(table_data)
        content.text_content.append({
            "type": "table",
            "content": table_data
        })
    
    # 提取图片
    if extract_images:
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_counter += 1
                    image_data = rel.target_part.blob
                    
                    # 确定图片格式
                    content_type = rel.target_part.content_type
                    if "png" in content_type:
                        ext = "png"
                    elif "jpeg" in content_type or "jpg" in content_type:
                        ext = "jpg"
                    elif "gif" in content_type:
                        ext = "gif"
                    else:
                        ext = "png"
                    
                    # 获取图片尺寸
                    width, height = 0, 0
                    if HAS_PIL:
                        try:
                            img = Image.open(io.BytesIO(image_data))
                            width, height = img.size
                        except:
                            pass
                    
                    extracted_img = ExtractedImage(
                        image_data=image_data,
                        image_ext=ext,
                        index=image_counter,
                        width=width,
                        height=height
                    )
                    
                    # 保存图片
                    img_filename = extracted_img.get_filename(file_path.stem)
                    img_path = images_dir / img_filename
                    img_path.write_bytes(image_data)
                    
                    extracted_img.image_data = b""  # 清空数据节省内存
                    content.images.append(extracted_img)
                except Exception as e:
                    continue
    
    content.total_images = image_counter
    return content


# ========== PPTX 解析 ==========

def extract_pptx_content(
    file_path: Path,
    output_dir: Path,
    images_subdir: str = "images",
    extract_images: bool = True
) -> OfficeContent:
    """
    提取PPTX文件内容
    """
    if not HAS_PPTX:
        raise ImportError("需要安装 python-pptx: pip install python-pptx")
    
    content = OfficeContent(file_type="pptx")
    prs = Presentation(str(file_path))
    
    # 提取元数据
    core_props = prs.core_properties
    content.metadata = {
        "title": core_props.title or "",
        "author": core_props.author or "",
        "subject": core_props.subject or "",
        "created": str(core_props.created) if core_props.created else "",
    }
    content.title = core_props.title or file_path.stem
    
    # 创建图片目录
    if extract_images:
        images_dir = output_dir / images_subdir
        images_dir.mkdir(parents=True, exist_ok=True)
    
    image_counter = 0
    
    # 遍历幻灯片
    for slide_num, slide in enumerate(prs.slides, 1):
        content.text_content.append({
            "type": "slide_marker",
            "content": f"--- 幻灯片 {slide_num} ---",
            "slide_num": slide_num
        })
        
        for shape in slide.shapes:
            # 提取文本
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    
                    # 根据字体大小判断标题级别
                    block_type = "paragraph"
                    if para.runs and para.runs[0].font.size:
                        font_size = para.runs[0].font.size.pt if para.runs[0].font.size else 12
                        if font_size >= 24:
                            block_type = "heading1"
                        elif font_size >= 18:
                            block_type = "heading2"
                        elif font_size >= 14:
                            block_type = "heading3"
                    
                    content.text_content.append({
                        "type": block_type,
                        "content": text,
                        "slide_num": slide_num
                    })
            
            # 提取表格
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                content.tables.append(table_data)
                content.text_content.append({
                    "type": "table",
                    "content": table_data,
                    "slide_num": slide_num
                })
            
            # 提取图片 - 使用hasattr检查而非硬编码shape_type
            if extract_images:
                # 检查是否为图片形状：优先使用 MSO_SHAPE_TYPE.PICTURE，否则用 hasattr
                is_picture = False
                if MSO_SHAPE_TYPE is not None:
                    try:
                        is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    except:
                        is_picture = hasattr(shape, 'image')
                else:
                    is_picture = hasattr(shape, 'image')
                
                if is_picture and hasattr(shape, 'image'):
                    try:
                        image_counter += 1
                        image = shape.image
                        image_data = image.blob
                        ext = image.ext
                        
                        width, height = 0, 0
                        if HAS_PIL:
                            try:
                                img = Image.open(io.BytesIO(image_data))
                                width, height = img.size
                            except:
                                pass
                        
                        extracted_img = ExtractedImage(
                            image_data=image_data,
                            image_ext=ext,
                            index=image_counter,
                            width=width,
                            height=height
                        )
                        
                        img_filename = extracted_img.get_filename(file_path.stem)
                        img_path = images_dir / img_filename
                        img_path.write_bytes(image_data)
                        
                        extracted_img.image_data = b""
                        content.images.append(extracted_img)
                        
                        content.text_content.append({
                            "type": "image",
                            "content": f"[图片 {image_counter}]",
                            "image_index": image_counter,
                            "image_ext": ext,
                            "slide_num": slide_num
                        })
                    except Exception as e:
                        continue
    
    content.total_images = image_counter
    return content


# ========== XLSX 解析 ==========

def extract_xlsx_content(
    file_path: Path,
    output_dir: Path,
    images_subdir: str = "images",
    extract_images: bool = True
) -> OfficeContent:
    """
    提取XLSX文件内容
    """
    if not HAS_OPENPYXL:
        raise ImportError("需要安装 openpyxl: pip install openpyxl")
    
    content = OfficeContent(file_type="xlsx")
    wb = load_workbook(str(file_path), data_only=True)
    
    # 提取元数据
    content.metadata = {
        "title": wb.properties.title or "",
        "author": wb.properties.creator or "",
        "created": str(wb.properties.created) if wb.properties.created else "",
    }
    content.title = wb.properties.title or file_path.stem
    
    # 创建图片目录
    if extract_images:
        images_dir = output_dir / images_subdir
        images_dir.mkdir(parents=True, exist_ok=True)
    
    image_counter = 0
    
    # 遍历工作表
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        
        content.text_content.append({
            "type": "sheet_marker",
            "content": f"## 工作表: {sheet_name}",
            "sheet_name": sheet_name
        })
        
        # 提取数据作为表格
        table_data = []
        for row in sheet.iter_rows(values_only=True):
            # 过滤全空行
            if all(cell is None for cell in row):
                continue
            row_data = [str(cell) if cell is not None else "" for cell in row]
            table_data.append(row_data)
        
        if table_data:
            content.tables.append(table_data)
            content.text_content.append({
                "type": "table",
                "content": table_data,
                "sheet_name": sheet_name
            })
        
        # 提取图片 - openpyxl通过 drawing.image 访问
        if extract_images:
            try:
                # openpyxl 3.0+ 使用 sheet._images 或遍历 drawing
                images_list = []
                
                # 方法1: 尝试 _images 属性 (openpyxl 内部)
                if hasattr(sheet, '_images') and sheet._images:
                    images_list = list(sheet._images)
                
                # 方法2: 尝试遍历 _drawing (更可靠)
                if not images_list and hasattr(sheet, '_drawing') and sheet._drawing:
                    for chart_or_image in sheet._drawing:
                        if hasattr(chart_or_image, '_data'):
                            images_list.append(chart_or_image)
                
                for img_obj in images_list:
                    try:
                        image_counter += 1
                        # 获取图片数据
                        if hasattr(img_obj, '_data'):
                            if callable(img_obj._data):
                                image_data = img_obj._data()
                            else:
                                image_data = img_obj._data
                        elif hasattr(img_obj, 'ref') and hasattr(img_obj.ref, 'blob'):
                            image_data = img_obj.ref.blob
                        else:
                            continue
                        
                        # 确定扩展名
                        ext = "png"
                        if hasattr(img_obj, 'format'):
                            ext = img_obj.format or "png"
                        
                        extracted_img = ExtractedImage(
                            image_data=image_data,
                            image_ext=ext,
                            index=image_counter
                        )
                        
                        img_filename = extracted_img.get_filename(file_path.stem)
                        img_path = images_dir / img_filename
                        img_path.write_bytes(image_data)
                        
                        extracted_img.image_data = b""
                        content.images.append(extracted_img)
                    except Exception as e:
                        image_counter -= 1  # 恢复计数器
                        continue
            except Exception as e:
                pass  # 图片提取失败不影响主流程
    
    wb.close()
    content.total_images = image_counter
    return content


# ========== RTF 解析 ==========

def extract_rtf_content(
    file_path: Path,
    output_dir: Path,
    images_subdir: str = "images",
    extract_images: bool = True
) -> OfficeContent:
    """
    提取RTF文件内容
    
    Args:
        file_path: RTF文件路径
        output_dir: 输出目录
        images_subdir: 图片子目录名
        extract_images: 是否提取图片（RTF中嵌入的图片支持有限）
    
    Returns:
        OfficeContent: 提取的内容
    """
    if not HAS_RTF:
        raise ImportError("需要安装 striprtf: pip install striprtf")
    
    content = OfficeContent(file_type="rtf")
    
    # 读取RTF文件
    try:
        # 尝试多种编码
        rtf_text = None
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    rtf_text = f.read()
                break
            except UnicodeDecodeError:
                continue
        
        if rtf_text is None:
            # 最后尝试二进制读取
            with open(file_path, 'rb') as f:
                rtf_bytes = f.read()
            rtf_text = rtf_bytes.decode('latin-1', errors='replace')
        
        # 使用 striprtf 提取纯文本
        plain_text = rtf_to_text(rtf_text)
        
    except Exception as e:
        raise ValueError(f"无法读取RTF文件: {e}")
    
    # 设置元数据
    content.metadata = {
        "title": file_path.stem,
        "author": "",
        "created": "",
    }
    content.title = file_path.stem
    
    # 解析文本内容
    lines = plain_text.split('\n')
    current_paragraph = []
    
    for line in lines:
        line = line.strip()
        
        if not line:
            # 空行表示段落结束
            if current_paragraph:
                para_text = ' '.join(current_paragraph)
                # 判断是否为标题（简单启发式：短行且不以标点结尾）
                if len(para_text) < 50 and not para_text.endswith(('.', '。', '!', '！', '?', '？', ',')):
                    content.text_content.append({
                        "type": "heading1",
                        "content": para_text
                    })
                else:
                    content.text_content.append({
                        "type": "paragraph",
                        "content": para_text
                    })
                current_paragraph = []
        else:
            # 检测列表项
            if line.startswith(('• ', '- ', '* ', '· ')):
                if current_paragraph:
                    content.text_content.append({
                        "type": "paragraph",
                        "content": ' '.join(current_paragraph)
                    })
                    current_paragraph = []
                content.text_content.append({
                    "type": "list_item",
                    "content": line[2:].strip()
                })
            elif len(line) > 2 and line[0].isdigit() and line[1] in '.）)':
                # 数字列表
                if current_paragraph:
                    content.text_content.append({
                        "type": "paragraph",
                        "content": ' '.join(current_paragraph)
                    })
                    current_paragraph = []
                content.text_content.append({
                    "type": "list_item",
                    "content": line[2:].strip()
                })
            else:
                current_paragraph.append(line)
    
    # 处理最后一个段落
    if current_paragraph:
        content.text_content.append({
            "type": "paragraph",
            "content": ' '.join(current_paragraph)
        })
    
    # RTF 中的嵌入图片处理（有限支持）
    # striprtf 不直接支持图片提取，这里只记录信息
    content.total_images = 0
    
    return content



def convert_old_format_to_new(file_path: Path, temp_dir: Path) -> Optional[Path]:
    """
    将旧格式Office文件转换为新格式
    
    Args:
        file_path: 旧格式文件路径 (.doc, .ppt, .xls)
        temp_dir: 临时目录
    
    Returns:
        转换后的新格式文件路径，失败返回None
    """
    # 延迟初始化 win32com
    if not _init_win32com():
        return None
    
    suffix = file_path.suffix.lower()
    format_map = {
        ".doc": ".docx",
        ".ppt": ".pptx", 
        ".xls": ".xlsx"
    }
    
    if suffix not in format_map:
        return None
    
    new_path = temp_dir / (file_path.stem + format_map[suffix])
    app = None
    doc = None
    
    try:
        pythoncom.CoInitialize()
        
        if suffix == ".doc":
            app = win32com.Dispatch("Word.Application")
            app.Visible = False
            app.DisplayAlerts = False  # 禁止弹窗
            try:
                doc = app.Documents.Open(
                    str(file_path.absolute()),
                    ReadOnly=True,
                    AddToRecentFiles=False
                )
                doc.SaveAs2(str(new_path.absolute()), FileFormat=16)  # docx format
            finally:
                if doc:
                    try:
                        doc.Close(SaveChanges=False)
                    except:
                        pass
                if app:
                    try:
                        app.Quit()
                    except:
                        pass
            
        elif suffix == ".ppt":
            app = win32com.Dispatch("PowerPoint.Application")
            # PowerPoint Visible 必须设为 True 或使用 msoFalse
            try:
                doc = app.Presentations.Open(
                    str(file_path.absolute()),
                    ReadOnly=True,
                    WithWindow=False
                )
                doc.SaveAs(str(new_path.absolute()), FileFormat=24)  # pptx format
            finally:
                if doc:
                    try:
                        doc.Close()
                    except:
                        pass
                if app:
                    try:
                        app.Quit()
                    except:
                        pass
            
        elif suffix == ".xls":
            app = win32com.Dispatch("Excel.Application")
            app.Visible = False
            app.DisplayAlerts = False  # 禁止弹窗
            try:
                doc = app.Workbooks.Open(
                    str(file_path.absolute()),
                    ReadOnly=True,
                    UpdateLinks=False
                )
                doc.SaveAs(str(new_path.absolute()), FileFormat=51)  # xlsx format
            finally:
                if doc:
                    try:
                        doc.Close(SaveChanges=False)
                    except:
                        pass
                if app:
                    try:
                        app.Quit()
                    except:
                        pass
        
        return new_path if new_path.exists() else None
        
    except Exception as e:
        # 确保清理COM对象
        if doc:
            try:
                doc.Close()
            except:
                pass
        if app:
            try:
                app.Quit()
            except:
                pass
        return None
    finally:
        try:
            pythoncom.CoUninitialize()
        except:
            pass


# ========== 通用接口 ==========

def extract_office_content(
    file_path: Path,
    output_dir: Path,
    images_subdir: str = "images",
    extract_images: bool = True
) -> OfficeContent:
    """
    提取Office文档内容（通用接口）
    
    Args:
        file_path: Office文档路径
        output_dir: 输出目录
        images_subdir: 图片子目录名
        extract_images: 是否提取图片
    
    Returns:
        OfficeContent: 提取的内容
    """
    suffix = file_path.suffix.lower()
    
    # 新格式直接解析
    if suffix == ".docx":
        return extract_docx_content(file_path, output_dir, images_subdir, extract_images)
    elif suffix == ".pptx":
        return extract_pptx_content(file_path, output_dir, images_subdir, extract_images)
    elif suffix == ".xlsx":
        return extract_xlsx_content(file_path, output_dir, images_subdir, extract_images)
    elif suffix == ".rtf":
        return extract_rtf_content(file_path, output_dir, images_subdir, extract_images)
    
    # 旧格式需要转换
    elif suffix in [".doc", ".ppt", ".xls"]:
        # 尝试初始化 win32com
        if not _init_win32com():
            raise ImportError(
                f"处理 {suffix} 格式需要安装 pywin32 并确保安装了对应的 Office 软件。\n"
                f"请运行: pip install pywin32\n"
                f"或者将文件另存为新格式 ({suffix}x)"
            )
        
        # 创建临时目录
        temp_dir = Path(tempfile.mkdtemp())
        try:
            new_path = convert_old_format_to_new(file_path, temp_dir)
            if not new_path:
                raise Exception(f"无法转换 {suffix} 格式，请确保安装了对应的 Office 软件")
            
            # 递归调用处理新格式
            content = extract_office_content(new_path, output_dir, images_subdir, extract_images)
            content.file_type = suffix[1:]  # 记录原始格式
            return content
        finally:
            # 清理临时文件
            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
            except:
                pass
    
    else:
        raise ValueError(f"不支持的文件格式: {suffix}")


def office_content_to_markdown(
    content: OfficeContent,
    file_path: Path,
    images_subdir: str = "images"
) -> str:
    """
    将Office内容转换为Markdown
    
    Args:
        content: Office文档内容
        file_path: 原文件路径
        images_subdir: 图片子目录
    
    Returns:
        Markdown文本
    """
    from urllib.parse import quote
    
    lines = []
    
    # 文档头
    title = content.title or file_path.stem
    lines.append(f"# {title}")
    lines.append("")
    
    # 元信息
    lines.append(f"> **源文件名**: {file_path.name}")
    lines.append(f"> **源文件绝对路径**: `{file_path.absolute()}`")
    lines.append(f"> **文件类型**: {content.file_type.upper()}")
    if content.metadata.get("author"):
        lines.append(f"> **作者**: {content.metadata['author']}")
    if content.metadata.get("created"):
        lines.append(f"> **创建时间**: {content.metadata['created']}")
    try:
        file_size = file_path.stat().st_size / 1024
        lines.append(f"> **文件大小**: {file_size:.1f} KB")
    except:
        pass
    lines.append(f"> **转换时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    if content.total_images > 0:
        lines.append(f"> **提取图片**: {content.total_images} 张")
    lines.append("")
    lines.append("---")
    lines.append("")
    
    # 转换内容
    image_counter = 0
    for block in content.text_content:
        block_type = block.get("type", "paragraph")
        block_content = block.get("content", "")
        
        if block_type == "heading1":
            lines.append(f"## {block_content}")
            lines.append("")
        elif block_type == "heading2":
            lines.append(f"### {block_content}")
            lines.append("")
        elif block_type == "heading3":
            lines.append(f"#### {block_content}")
            lines.append("")
        elif block_type == "list_item":
            lines.append(f"- {block_content}")
        elif block_type == "slide_marker":
            lines.append("")
            lines.append(f"---")
            lines.append("")
            lines.append(f"### {block_content}")
            lines.append("")
        elif block_type == "sheet_marker":
            lines.append("")
            lines.append(f"{block_content}")
            lines.append("")
        elif block_type == "table":
            table_data = block_content
            if isinstance(table_data, list) and table_data:
                # 生成Markdown表格
                lines.append("")
                # 表头
                header = table_data[0]
                lines.append("| " + " | ".join(str(cell) for cell in header) + " |")
                lines.append("| " + " | ".join("---" for _ in header) + " |")
                # 数据行
                for row in table_data[1:]:
                    lines.append("| " + " | ".join(str(cell) for cell in row) + " |")
                lines.append("")
        elif block_type == "image":
            image_counter += 1
            img_index = block.get("image_index", image_counter)
            # 从 content.images 获取正确的文件名（包含实际扩展名）
            img_filename = None
            for img in content.images:
                if img.index == img_index:
                    img_filename = img.get_filename(file_path.stem)
                    break
            # 如果找不到，使用 block 中的扩展名，否则默认 png
            if not img_filename:
                img_ext = block.get("image_ext", "png")
                img_filename = f"{file_path.stem}_img{img_index}.{img_ext}"
            img_path = f"{images_subdir}/{quote(img_filename)}"
            lines.append(f"![图片{img_index}]({img_path})")
            lines.append("")
        else:
            # 普通段落
            lines.append(block_content)
            lines.append("")
    
    # 添加未在 text_content 中引用的图片（DOCX/XLSX 的图片）
    referenced_indices = {
        block.get("image_index") 
        for block in content.text_content 
        if block.get("type") == "image"
    }
    unreferenced_images = [img for img in content.images if img.index not in referenced_indices]
    
    if unreferenced_images:
        lines.append("")
        lines.append("---")
        lines.append("")
        lines.append("## 文档图片")
        lines.append("")
        for img in unreferenced_images:
            img_filename = img.get_filename(file_path.stem)
            img_path = f"{images_subdir}/{quote(img_filename)}"
            alt_text = f"图片{img.index}"
            if img.width and img.height:
                alt_text += f" ({img.width}x{img.height})"
            lines.append(f"![{alt_text}]({img_path})")
            lines.append("")
    
    return "\n".join(lines)


# 检查依赖状态
def check_dependencies() -> Dict[str, bool]:
    """检查Office解析依赖状态"""
    return {
        "python-docx": HAS_DOCX,
        "python-pptx": HAS_PPTX,
        "openpyxl": HAS_OPENPYXL,
        "pywin32": HAS_WIN32COM,
        "PIL": HAS_PIL,
        "striprtf": HAS_RTF,
    }


# 获取支持的文件扩展名
def get_supported_extensions() -> List[str]:
    """获取支持的文件扩展名"""
    extensions = []
    
    if HAS_DOCX:
        extensions.extend([".docx"])
    if HAS_PPTX:
        extensions.extend([".pptx"])
    if HAS_OPENPYXL:
        extensions.extend([".xlsx"])
    if HAS_RTF:
        extensions.append(".rtf")
    
    # 旧格式需要pywin32
    if HAS_WIN32COM:
        if HAS_DOCX:
            extensions.append(".doc")
        if HAS_PPTX:
            extensions.append(".ppt")
        if HAS_OPENPYXL:
            extensions.append(".xls")
    
    return extensions
